import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import random
from PIL import Image
import numpy as np
from collections import Counter
from concurrent.futures import ThreadPoolExecutor




def add_slide_with_image(img_path, color):
    # 添加一页幻灯片
    slide_layout = presentation.slide_layouts[6]  # 选择一种布局（6: 空白布局）
    slide = presentation.slides.add_slide(slide_layout)

    # 定义幻灯片的尺寸（16:9 比例）
    slide_width = presentation.slide_width
    slide_height = presentation.slide_height

    # 图片的原始尺寸
    original_width = 1920
    original_height = 882

    # 计算纵横比
    aspect_ratio = original_width / original_height

    # 计算调整后的宽度和高度
    if (slide_width / slide_height) > aspect_ratio:
        # 以高度为基准
        height = slide_height
        width = height * aspect_ratio
    else:
        # 以宽度为基准
        width = slide_width
        height = width / aspect_ratio

    # 计算图片位置（居中）
    left = (slide_width - width) / 2
    top = (slide_height - height) / 2

    # 添加图片
    slide.shapes.add_picture(img_path, left, top, width=width, height=height)

    # 添加文本框
    text_left = left
    text_top = top + height
    text_width = width
    text_height = Inches(1)

    textbox = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
    text_frame = textbox.text_frame

    # 添加艺术字文本
    p = text_frame.add_paragraph()
    p.text = os.path.basename(img_path)[:-4]  # 去掉文件后缀

    # 设置字体样式
    run = p.runs[0]
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = color
    run.font.name = 'Arial'

    # 设置文本居中
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER

def get_dominant_color(image_path):
    # 打开图像
    image = Image.open(image_path)
    # 将图像转换为 RGB 模式
    image = image.convert('RGB')

    # 将图像数据转换为 numpy 数组
    pixels = np.array(image)
    # 将二维数组展平为一维
    pixels = pixels.reshape(-1, 3)

    # 计算颜色出现的频率
    color_counts = Counter(map(tuple, pixels))

    # 找到出现频率最高的颜色
    dominant_color = color_counts.most_common(1)[0][0]
    # 将颜色转换为 RGBColor 对象
    # 确保颜色值在 0-255 范围内
    r = int(max(0, min(255, dominant_color[0])))
    g = int(max(0, min(255, dominant_color[1])))
    b = int(max(0, min(255, dominant_color[2])))

    # 将颜色转换为 RGBColor 对象
    dominant_color_rgb = RGBColor(r, g, b)

    return dominant_color_rgb




# 创建一个新的 PowerPoint 演示文稿
presentation = Presentation()



# 记录运行时间
import time
# 记录开始时间
start_time = time.time()
# 遍历该文件夹下的文件夹
folderPathHead = "./hero_name"
# with ThreadPoolExecutor() as executor:
for folder in os.listdir(folderPathHead):
    folderPath = os.path.join(folderPathHead, folder)
    # 随机颜色
    color = RGBColor(random.randint(0, 255), random.randint(0, 255), random.randint(0, 255))
    for filename in os.listdir(folderPath):
        if filename.endswith(".jpg"):
            img_path = os.path.join(folderPath, filename)
            # color = get_dominant_color(img_path)
            add_slide_with_image(img_path, color)
            # executor.submit(add_slide_with_image, img_path,color)



# 保存演示文稿
presentation.save('output.pptx')  # 输出文件名
print("演示文稿已生成！")
# 记录结束时间
end_time = time.time()
# 计算并打印总耗时
total_time = end_time - start_time
print(f"下载完成，总耗时: {total_time:.2f} 秒")
